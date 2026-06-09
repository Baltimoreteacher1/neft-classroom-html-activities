#!/usr/bin/env python3
"""extract_reveal.py — pull curated teaching pieces out of a Reveal Math .pptx.

Given a Reveal Math lesson deck, this finds two curated pieces the HTML lesson
engine knows how to render:

  1. Notice & Wonder  — the Be Curious LAUNCH IMAGE students notice & wonder
     about (the engaging real-world photo/graphic on the Be Curious slide),
     plus a short scenario/context sentence.
  2. Word problem     — the application/"Apply" slide: its image (if any), the
     problem text, and a short title from the "Apply: ___" label.

It is a *read-only reporter*: it extracts text, copies the chosen raster images
to the requested output paths (web-sized; format matched to source), and prints
a single JSON blob on stdout. All config wiring + git/build orchestration lives
in scripts/reveal-lesson.mjs.

Usage:
  python3 extract_reveal.py <deck.pptx> --notice-wonder-out <path> \
      --word-problem-out <path> [--no-write]   # the extension is re-derived

With --no-write, no image files are written (dry-run); the JSON still reports
which slides/images WOULD be used and their dimensions.

Heuristics (robust across the ~70 Grade-6 Reveal decks):

  * Notice & Wonder image: anchor on the "Be Curious / Mindset" slide (the
    launch routine varies — "What do you notice/wonder?" vs "What could the
    question be?"), then take the LARGEST real raster image on that slide (or
    the next slide that has one). Images inside group shapes and picture
    placeholders are included; vector (EMF/WMF/SVG) and undecodable images are
    skipped so a broken asset is never written.
  * Notice & Wonder context: the longest descriptive scenario block in a small
    window after the anchor, with procedural rule/summary text excluded.
  * Word problem: the "Apply" slide (Apply as a heading/title, not a substring),
    weighted toward late-deck slides; its largest real image; title from the
    "Apply: ___" label.
"""

from __future__ import annotations

import argparse
import io
import json
import os
import re
import sys

try:
    from PIL import Image

    # Guard untrusted decks against decompression-bomb images.
    Image.MAX_IMAGE_PIXELS = 64_000_000  # ~64 MP
except Exception:  # Pillow optional — without it we can't size/re-encode images.
    Image = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
except Exception as exc:  # pragma: no cover - dependency guard
    print(json.dumps({"error": f"python-pptx import failed: {exc}"}), file=sys.stdout)
    sys.exit(2)


# Raster formats a browser can display (and Pillow can re-encode).
_RASTER_TYPES = ("image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp")

# Cap extracted text so a malformed/bloated deck can't write a huge config value.
_MAX_CONTEXT_CHARS = 600


def _clean_text(text: str) -> str:
    """Collapse whitespace and strip Reveal's invisible markers."""
    if not text:
        return ""
    # Drop zero-width / soft markers Reveal embeds, normalise unicode spaces.
    text = text.replace("​", "").replace("\xa0", " ")
    text = re.sub(r"[ \t]*\n[ \t]*", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _slide_text_blocks(slide):
    """Return cleaned, non-empty text frames for a slide (order preserved)."""
    blocks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            raw = shape.text_frame.text
            cleaned = _clean_text(raw)
            if cleaned:
                blocks.append(cleaned)
    return blocks


def _iter_picture_shapes(shapes):
    """Yield leaf picture shapes, recursing into groups. Reveal commonly groups
    the launch photo / data graphic with captions, and uses picture placeholders
    — both are invisible to a flat ``shape_type == PICTURE`` scan."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
            continue
        # Picture placeholder: carries .image but isn't a PICTURE shape_type.
        try:
            if shape.image is not None:
                yield shape
        except Exception:
            continue


def _slide_pictures(slide):
    """Return raster picture dicts for a slide (groups + placeholders included).
    Non-raster (EMF/WMF/SVG) and undecodable images are dropped so we never write
    a browser-unreadable asset under a .png/.jpg name."""
    pics = []
    for shape in _iter_picture_shapes(slide.shapes):
        try:
            image = shape.image
            blob = image.blob
            content_type = (image.content_type or "").lower()
            partname = getattr(image, "filename", None) or ""
        except Exception:
            continue
        if not any(t in content_type for t in _RASTER_TYPES):
            continue  # vector/unsupported — a browser can't show it
        if Image is None:
            continue  # can't verify/encode without Pillow
        try:
            with Image.open(io.BytesIO(blob)) as im:
                w, h = im.size
        except Exception:
            continue  # corrupt/undecodable — skip rather than ship broken art
        emu_w = Emu(shape.width).inches if getattr(shape, "width", None) else None
        emu_h = Emu(shape.height).inches if getattr(shape, "height", None) else None
        pics.append(
            {
                "partname": partname,
                "content_type": content_type,
                "bytes": len(blob),
                "w": w,
                "h": h,
                "emu_w": emu_w,
                "emu_h": emu_h,
                "blob": blob,
            }
        )
    return pics


# Banners/badges are tiny; a real launch/data image clears this floor.
_MIN_IMAGE_PX = 40_000


def _is_real_image(pic) -> bool:
    """True for a substantial raster image, not a tiny banner/icon. Unknown
    dimensions are rejected (every kept pic was decoded in _slide_pictures)."""
    w, h = pic.get("w"), pic.get("h")
    return bool(w and h and (w * h) >= _MIN_IMAGE_PX)


def _image_area(pic) -> float:
    return float((pic.get("w") or 0) * (pic.get("h") or 0))


def _find_notice_wonder(slides):
    """Locate the Be Curious slide, its PHOTO (the launch image students notice &
    wonder about), and the data-context text.

    Teacher preference: the Notice & Wonder image is the Be Curious launch PHOTO
    (the engaging real-world image on the "What do you notice? / wonder?" slide),
    NOT the later data table. We pick the largest real image on the Be Curious
    anchor slide (or the next slide that has one).
    """
    # 1) Anchor on the "Be Curious / notice / wonder / mindset" slides.
    anchors = []
    for idx, slide in enumerate(slides):
        joined = " ".join(_slide_text_blocks(slide)).lower()
        # Reveal's Be Curious launch varies: "What do you notice? / wonder?" on
        # some decks, "What could the question be?" on others. Anchor on the
        # Be Curious / Mindset marker itself (don't require notice/wonder words).
        if "be curious" in joined or "mindset" in joined:
            anchors.append(idx)

    # 2) The Be Curious photo: largest real image on the anchor slide, scanning a
    #    small window forward if the anchor itself has no usable image.
    if anchors:
        order = list(range(anchors[0], min(anchors[0] + 4, len(slides))))
    else:
        order = list(range(min(8, len(slides))))

    best_idx, best_pic = None, None
    for idx in order:
        pics = [p for p in _slide_pictures(slides[idx]) if _is_real_image(p)]
        if pics:
            best_idx = idx
            best_pic = max(pics, key=_image_area)
            break

    if best_pic is None:
        return None

    # 3) Context = the scenario/data sentence NEAR the Be Curious slide (not the
    #    later rule/summary). Gather text from the anchor slide + the next few,
    #    then pick the best descriptive (non-question, non-rule) block.
    start = anchors[0] if anchors else best_idx
    window_blocks = []
    for idx in range(start, min(start + 4, len(slides))):
        window_blocks.extend(_slide_text_blocks(slides[idx]))
    context = _pick_context(window_blocks)
    if not context:
        context = _pick_context(_slide_text_blocks(slides[best_idx]))

    return {
        "slide_number": best_idx + 1,
        "context": context,
        "image": best_pic,
    }


# Procedural rule/summary phrasing — NOT a notice/wonder scenario.
_RULE_HINTS = (
    "algorithm",
    "place the decimal",
    "multiply both",
    "power of 10",
    "quotient above",
    "you can use",
)


def _pick_context(blocks):
    """Pick the scenario/data sentence from a set of text blocks."""
    descriptive = []
    for b in blocks:
        low = b.lower()
        if low in ("reveal:", "think about it:", "be curious", "mindset"):
            continue
        if low.endswith(":") and len(b) < 40:
            continue  # short labels like "Teaching Experience"
        if any(r in low for r in _RULE_HINTS):
            continue  # skip the procedural rule/summary
        descriptive.append(b)
    if not descriptive:
        return ""
    # The richest scenario/data block is reliably the longest one (it may end
    # with the framing question, which is fine for a Notice & Wonder prompt).
    # Only fall back to skipping questions if the longest is a bare mindset ask.
    longest = max(descriptive, key=len)
    if longest.strip().endswith("?") and len(longest) < 30:
        non_q = [b for b in descriptive if not b.strip().endswith("?")]
        if non_q:
            longest = max(non_q, key=len)
    return longest[:_MAX_CONTEXT_CHARS].strip()


def _find_word_problem(slides):
    """Locate the application / "Apply:" slide, its image, text, and title."""
    candidates = []
    n = len(slides)
    for idx, slide in enumerate(slides):
        blocks = _slide_text_blocks(slide)
        # Require "Apply" as a slide heading (Reveal's Apply section), not a
        # substring like "apply the formula" on a teaching slide.
        has_apply_label = any(b.lower().startswith("apply") for b in blocks)
        if not has_apply_label:
            continue
        # A real application problem asks a question; weight late-deck position.
        has_question = "?" in " ".join(blocks)
        lateness = idx / max(n - 1, 1)
        score = (1.0 if has_question else 0.0) + lateness
        candidates.append((score, idx, slide, blocks))

    if not candidates:
        return None

    candidates.sort(key=lambda c: c[0], reverse=True)
    _, idx, slide, blocks = candidates[0]

    title = _derive_apply_title(blocks)
    text = _pick_problem_text(blocks)

    pics = [p for p in _slide_pictures(slide) if _is_real_image(p)]
    image = max(pics, key=_image_area) if pics else None

    return {
        "slide_number": idx + 1,
        "title": title,
        "text": text,
        "image": image,
    }


def _derive_apply_title(blocks):
    """Turn an "Apply: Real Estate" label into a clean title."""
    for b in blocks:
        if b.lower().startswith("apply"):
            # "Apply: Real Estate" -> "Apply: Real Estate"
            label = re.sub(r"\s+", " ", b).strip()
            return label
    return "Apply"


def _pick_problem_text(blocks):
    """Pick the real-world problem statement from the slide."""
    body = [b for b in blocks if not b.lower().startswith("apply")]
    if not body:
        return ""
    text = max(body, key=len)
    # Reveal often prefixes the ask with a literal "Question:" — keep the
    # whole statement but normalise that marker into a sentence flow.
    text = re.sub(r"\bQuestion:\s*", "", text).strip()
    return text[:_MAX_CONTEXT_CHARS].strip()


def _web_extension(content_type) -> str:
    """Keep photos as .jpg (light) and graphics as .png."""
    ct = (content_type or "").lower()
    return ".jpg" if ("jpeg" in ct or "jpg" in ct) else ".png"


def _save_web_image(blob, final_path, ext, max_w=1100):
    """Downscale to <= max_w wide and save in a web-friendly format."""
    if Image is None:
        with open(final_path, "wb") as fh:
            fh.write(blob)
        return
    try:
        img = Image.open(io.BytesIO(blob))
        if img.width and img.width > max_w:
            ratio = max_w / float(img.width)
            img = img.resize((max_w, max(1, int(img.height * ratio))), Image.LANCZOS)
        if ext == ".jpg":
            img.convert("RGB").save(final_path, "JPEG", quality=85, optimize=True)
        else:
            img.save(final_path, "PNG", optimize=True)
    except Exception:
        with open(final_path, "wb") as fh:
            fh.write(blob)


def _emit_image(pic, out_path, write):
    """Write the chosen image (web-sized, format matched to source); return a
    report whose `out`/`filename` reflect the ACTUAL extension written."""
    if pic is None:
        return None
    ext = _web_extension(pic.get("content_type"))
    base, _ = os.path.splitext(out_path)
    final_path = base + ext
    report = {
        "partname": pic["partname"],
        "content_type": pic["content_type"],
        "bytes": pic["bytes"],
        "width": pic["w"],
        "height": pic["h"],
        "out": final_path,
        "filename": os.path.basename(final_path),
    }
    if write:
        _save_web_image(pic["blob"], final_path, ext)
    return report


def main(argv=None):
    ap = argparse.ArgumentParser(description="Extract curated Reveal pieces.")
    ap.add_argument("pptx")
    ap.add_argument("--notice-wonder-out", required=True)
    ap.add_argument("--word-problem-out", required=True)
    ap.add_argument(
        "--no-write",
        action="store_true",
        help="Report only; do not write any image files.",
    )
    args = ap.parse_args(argv)

    try:
        prs = Presentation(args.pptx)
    except Exception as exc:
        print(json.dumps({"error": f"could not open pptx: {exc}"}))
        return 2

    slides = list(prs.slides)
    write = not args.no_write

    nw = _find_notice_wonder(slides)
    wp = _find_word_problem(slides)

    out = {"slide_count": len(slides), "noticeAndWonder": None, "wordProblem": None}

    if nw is not None:
        img_report = _emit_image(nw["image"], args.notice_wonder_out, write)
        out["noticeAndWonder"] = {
            "slide_number": nw["slide_number"],
            "context": nw["context"],
            "image": img_report,
        }

    if wp is not None:
        img_report = _emit_image(wp["image"], args.word_problem_out, write)
        out["wordProblem"] = {
            "slide_number": wp["slide_number"],
            "title": wp["title"],
            "text": wp["text"],
            "image": img_report,
        }

    print(json.dumps(out))
    return 0


if __name__ == "__main__":
    sys.exit(main())
